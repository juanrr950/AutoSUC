
from pptx import Presentation
from pptx.util import Pt
import os
from AutoSUC.settings import BASE_DIR
def generar_powerpoint(suc):
        
    # Create an instance of Presentation class
    def find_shape(slide,name):
        
        for i in slide.shapes:
            if name==i.name:
                return i
    def delete_shape(name_shape):  
        try:
            elem=find_shape(sld, name_shape)
            elem._element.getparent().remove(elem._element)
        except:
            pass
    n=suc.num_postes
    
    if n<=5:
        sufijo="_5"
    elif n<=7:
        sufijo="_7"
    else:
        raise NotImplementedError("PLANTILLAS POWER POINT NO IMPLEMENTADAS PARA MAS DE 7 POSTES")
    
    if suc.plantilla=='H':
        pres = Presentation(os.path.join(BASE_DIR,
                            'media/plantillas/PLANTILLA_HORIZONTAL'+sufijo+'.pptx'))
    else:
        pres = Presentation(os.path.join(BASE_DIR,
                            'media/plantillas/PLANTILLA_VERTICAL'+sufijo+'.pptx'))
    if n>7:
        raise NotImplementedError("PLANTILLAS POWER POINT NO IMPLEMENTADAS PARA MAS DE 7 POSTES")
    # Get the first slide
    sld = pres.slides[0]
    
    #Cambiamos los postes 
    p1=find_shape(sld, "p1")
    p1.text_frame.paragraphs[0].runs[0].text=suc.poste_1
    
    p2=find_shape(sld, "p2")
    p2.text_frame.paragraphs[0].runs[0].text=suc.poste_2
        
    
    #Introducimos las medidas
    m1=find_shape(sld, "m1")
    m1.text_frame.paragraphs[0].runs[0].text=suc.medida_1_2 + " m."
       
    n=suc.num_postes
    if n>=3:
        p3=find_shape(sld, "p3")
        p3.text_frame.paragraphs[0].runs[0].text=suc.poste_3
        
        m2=find_shape(sld, "m2")
        m2.text_frame.paragraphs[0].runs[0].text=suc.medida_2_3 + " m."
    if n>=4:
        p4=find_shape(sld, "p4")
        p4.text_frame.paragraphs[0].runs[0].text=suc.poste_4
    
        m3=find_shape(sld, "m3")
        m3.text_frame.paragraphs[0].runs[0].text=suc.medida_3_4 + " m."
    if n>=5:
        p5=find_shape(sld, "p5")
        p5.text_frame.paragraphs[0].runs[0].text=suc.poste_5
        
        m4=find_shape(sld, "m4")
        m4.text_frame.paragraphs[0].runs[0].text=suc.medida_4_5 + " m."
    
    if n>=6:
        p6=find_shape(sld, "p6")
        p6.text_frame.paragraphs[0].runs[0].text=suc.poste_6
        
        m5=find_shape(sld, "m5")
        m5.text_frame.paragraphs[0].runs[0].text=suc.medida_5_6 + " m."
         
    if n>=7:
        p7=find_shape(sld, "p7")
        p7.text_frame.paragraphs[0].runs[0].text=suc.poste_7
        
        m6=find_shape(sld, "m6")
        m6.text_frame.paragraphs[0].runs[0].text=suc.medida_6_7 + " m."     
    
    if n>=8:
        p8=find_shape(sld, "p8")
        p8.text_frame.paragraphs[0].runs[0].text=suc.poste_8
        
        m7=find_shape(sld, "m7")
        m7.text_frame.paragraphs[0].runs[0].text=suc.medida_7_8 + " m."     
        
    if n>=9:
        p9=find_shape(sld, "p9")
        p9.text_frame.paragraphs[0].runs[0].text=suc.poste_9
        
        m8=find_shape(sld, "m8")
        m8.text_frame.paragraphs[0].runs[0].text=suc.medida_8_9 + " m."      
    
    if n>=10:
        p10=find_shape(sld, "p10")
        p10.text_frame.paragraphs[0].runs[0].text=suc.poste_10
        
        m9=find_shape(sld, "m9")
        m9.text_frame.paragraphs[0].runs[0].text=suc.medida_9_10 + " m."    
          
    #Borramos en cada caso si son menos de 5 postes
    if n<=9:    
        delete_shape("p10")
        delete_shape("vortex10")
        delete_shape("conector8_10")
        delete_shape("m9")
        delete_shape("letter9_10") 
    
    if n<=8:    
        delete_shape("p9")
        delete_shape("vortex9")
        delete_shape("conector8_9")
        delete_shape("m8")
        delete_shape("letter8_9") 
    
    if n<=7:    
        delete_shape("p8")
        delete_shape("vortex8")
        delete_shape("conector6_8")
        delete_shape("m7")
        delete_shape("letter7_8")  
    
    if n<=6:    
        delete_shape("p7")
        delete_shape("vortex7")
        delete_shape("conector6_7")
        delete_shape("m6")
        delete_shape("letter6_7")  
    
    if n<=5:    
        delete_shape("p6")
        delete_shape("vortex6")
        delete_shape("conector5_6")
        delete_shape("m5")
        delete_shape("letter5_6")  
    
    if n<=4:    
        delete_shape("p5")
        delete_shape("vortex5")
        delete_shape("conector4_5")
        delete_shape("m4")
        delete_shape("letter4_5")       
       
    if n<=3:    
        delete_shape("p4")
        delete_shape("vortex4")
        delete_shape("conector3_4")
        delete_shape("m3")
        delete_shape("letter3_4")       
    if n<=2:    
        delete_shape("p3")
        delete_shape("vortex3")
        delete_shape("conector2_3")
        delete_shape("m2")
        delete_shape("letter2_3")   
    
    ct=find_shape(sld, "cartelct")
    ct.text_frame.paragraphs[0].runs[0].text=suc.nombre_central
    ct.text_frame.paragraphs[1].runs[0].text="MIGA: "+str(suc.codigo_miga)

    
    nombreppt=suc.nombre.split('_')[1]+suc.nombre.split('_')[2]
    nombreppt=''.join(e for e in nombreppt if e.isalnum())
    nombreppt=nombreppt[-11:]
    pres.save(os.path.join(BASE_DIR,
                        'media/'+suc.usuario.username+'/'+
                        suc.nombre+'/'+nombreppt+'.ppt'))
    
    suc.powerpoint=suc.usuario.username+'/'+suc.nombre+'/'+nombreppt+'.ppt'
    suc.save()